from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_std_presentation_with_images():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Sexually Transmitted Diseases: Comprehensive Overview"
    subtitle.text = "MBBS 3rd Year Class\nDuration: 60 minutes\nAuthor: Dr Siddalingaiah H S, Professor, Community Medicine, SIMSRH, Tumkur\nEmail: hssling@yahoo.com | Phone: 8941087719\nDate: [Date]"

    # Slide 2: What are STDs?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'What are STDs?'
    tf = body_shape.text_frame
    tf.text = 'Definition and Classification'

    p = tf.add_paragraph()
    p.text = 'Definition: Infections transmitted through sexual contact, including vaginal, anal, and oral sex'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Classification by Causative Organism:'
    p.level = 1

    # Add STD classification image
    try:
        left = Inches(2)
        top = Inches(2.5)
        height = Inches(4)
        slide.shapes.add_picture('std_classification.png', left, top, height=height)
    except:
        print("Could not add STD classification image")

    # Slide 3: Epidemiology - Global Burden
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Epidemiology - Global Burden'
    tf = body_shape.text_frame
    tf.text = 'WHO Statistics (2022): 1 million new STD cases daily'

    # Add global STD epidemiology chart
    try:
        left = Inches(1)
        top = Inches(2)
        height = Inches(5)
        slide.shapes.add_picture('global_std_epidemiology.png', left, top, height=height)
    except:
        print("Could not add global STD epidemiology chart")

    # Slide 4: Epidemiology - Indian Context
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Epidemiology - Indian Context'
    tf = body_shape.text_frame
    tf.text = 'NACO Estimates (2023):'

    p = tf.add_paragraph()
    p.text = '• 30-40 million STD cases annually'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HIV prevalence: 0.22% (23.1 lakh PLHIV)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Syphilis: Rising trend, especially congenital'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Gonorrhea/Chlamydia: High among youth and high-risk groups'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Regional Distribution:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Highest prevalence: Northeast (Nagaland, Manipur)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Southern states: Karnataka, Andhra Pradesh, Telangana'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Urban vs Rural: Higher in urban areas'
    p.level = 2

    # Slide 5: Transmission Routes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Transmission Routes'
    tf = body_shape.text_frame
    tf.text = 'How STDs Spread'

    p = tf.add_paragraph()
    p.text = 'Sexual Transmission:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Vaginal intercourse'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Anal intercourse'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Oral-genital contact'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Manual-genital contact'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Non-Sexual Transmission:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Mother-to-child (congenital syphilis, HIV)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Blood transfusion (HIV, Hepatitis B)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Sharing needles (HIV, Hepatitis B)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Organ transplantation'
    p.level = 2

    # Slide 6: Bacterial STDs - Gonorrhea
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Bacterial STDs - Gonorrhea'
    tf = body_shape.text_frame
    tf.text = 'Gonorrhea: Clinical Features'

    p = tf.add_paragraph()
    p.text = 'Causative Agent: Neisseria gonorrhoeae (Gram-negative diplococcus)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Presentation:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Males: Acute urethritis'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Purulent discharge (yellow/green)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Dysuria, frequency, urgency'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Incubation: 2-7 days'
    p.level = 3

    p = tf.add_paragraph()
    p.text = 'Females: Often asymptomatic (50%)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• May cause cervicitis, PID'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Abdominal pain, fever'
    p.level = 3

    # Slide 7: Bacterial STDs - Syphilis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Bacterial STDs - Syphilis'
    tf = body_shape.text_frame
    tf.text = 'Syphilis: The Great Imitator'

    p = tf.add_paragraph()
    p.text = 'Causative Agent: Treponema pallidum'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Stages:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '1. Primary (2-12 weeks):'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Painless chancre at inoculation site'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Clean base, raised borders'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Regional lymphadenopathy'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '2. Secondary (6-24 weeks):'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Generalized rash (palms/soles)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Condylomata lata, alopecia'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Fever, malaise, lymphadenopathy'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '3. Tertiary (>2 years):'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Cardiovascular syphilis'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Neurosyphilis, gummas'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Tabes dorsalis'
    p.level = 3

    # Slide 8: Bacterial STDs - Chlamydia
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Bacterial STDs - Chlamydia'
    tf = body_shape.text_frame
    tf.text = 'Chlamydia: Silent Infection'

    p = tf.add_paragraph()
    p.text = 'Causative Agent: Chlamydia trachomatis'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Features:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Often asymptomatic (70-80%)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Females: Cervicitis, PID, infertility'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Males: Urethritis, epididymitis'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Both: Proctitis, conjunctivitis'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Complications:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Pelvic inflammatory disease'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Ectopic pregnancy'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Chronic pelvic pain'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Infertility in both genders'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Key Fact: Most common bacterial STD worldwide'
    p.level = 1

    # Slide 9: Viral STDs Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Viral STDs Overview'
    tf = body_shape.text_frame
    tf.text = 'Viral STDs: Chronic Infections'

    p = tf.add_paragraph()
    p.text = 'HIV/AIDS:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Retrovirus attacking CD4+ T cells'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Progressive immunosuppression'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Lifelong infection, manageable with ART'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'HSV (Herpes Simplex):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HSV-1: Oral herpes'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• HSV-2: Genital herpes'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Recurrent painful ulcers'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Lifelong latency'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'HPV (Human Papillomavirus):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 100+ subtypes'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Low-risk: Genital warts'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• High-risk: Cervical cancer'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Vaccination available'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Hepatitis B:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Chronic liver disease'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Vaccine-preventable'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• High prevalence in India'
    p.level = 2

    # Slide 10: Diagnosis - Clinical Approach
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Diagnosis - Clinical Approach'
    tf = body_shape.text_frame
    tf.text = 'Diagnostic Strategy'

    p = tf.add_paragraph()
    p.text = '1. Sexual History (5 Ps):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Partners (number, type)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Practices (vaginal, anal, oral)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Protection (condom use)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Past STDs'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Pregnancy intentions'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '2. Physical Examination:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Genital inspection'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Lymph node palpation'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Systemic signs'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '3. Laboratory Tests:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Microscopy (Gram stain, wet mount)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Culture and sensitivity'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• PCR (gold standard for chlamydia)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Serology (syphilis, HIV)'
    p.level = 2

    # Slide 11: Diagnosis - Laboratory Tests
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Diagnosis - Laboratory Tests'
    tf = body_shape.text_frame
    tf.text = 'Specific Diagnostic Tests'

    p = tf.add_paragraph()
    p.text = 'Gonorrhea:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Gram stain: Intracellular diplococci'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Culture: Thayer-Martin medium'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• PCR: Most sensitive'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Syphilis:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• VDRL/TPHA: Screening'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• FTA-ABS: Confirmatory'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Dark field microscopy'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Chlamydia:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• PCR: Endocervical swab'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Culture: McCoy cells'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• EIA: Less sensitive'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'HSV:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• PCR: Vesicular fluid'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Viral culture'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Tzanck smear'
    p.level = 2

    # Slide 12: Treatment - General Principles
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Treatment - General Principles'
    tf = body_shape.text_frame
    tf.text = 'Treatment Guidelines'

    p = tf.add_paragraph()
    p.text = 'NACO STI Management Guidelines (2020):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Syndromic management approach'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Dual therapy for gonorrhea'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Partner treatment essential'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Test of cure recommended'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Key Principles:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Treat empirically based on symptoms'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Culture sensitivity for resistance'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Follow-up testing'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Prevention of reinfection'
    p.level = 2

    # Slide 13: Treatment - Specific Regimens
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Treatment - Specific Regimens'
    tf = body_shape.text_frame
    tf.text = 'Treatment Protocols'

    p = tf.add_paragraph()
    p.text = 'Syphilis:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Primary/Secondary: Benzathine penicillin 2.4 MU IM single dose'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Latent: Benzathine penicillin 2.4 MU IM weekly × 3'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Tertiary/Neurosyphilis: Aqueous penicillin G 3-4 MU IV q4h × 14 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Alternative: Doxycycline 100mg PO twice daily × 14 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Gonorrhea:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Ceftriaxone 500mg IM single dose'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• PLUS Azithromycin 1g PO single dose'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Test of cure in 7-14 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Chlamydia:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Azithromycin 1g PO single dose'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• OR Doxycycline 100mg PO twice daily × 7 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Test of cure recommended'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'HSV:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Acyclovir 400mg PO three times daily × 5-10 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Valacyclovir 1g PO twice daily × 5-10 days'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Suppressive therapy for recurrences'
    p.level = 2

    # Slide 14: Prevention and Control - Primary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention and Control - Primary'
    tf = body_shape.text_frame
    tf.text = 'Primary Prevention'

    p = tf.add_paragraph()
    p.text = 'ABC Approach:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Abstain from sex'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Be faithful to one partner'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Condoms consistently and correctly'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Vaccines:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HPV vaccine (9-26 years)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Hepatitis B vaccine'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• HIV vaccine (in development)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Other Strategies:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Pre-exposure prophylaxis (PrEP) for HIV'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Post-exposure prophylaxis (PEP)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Male circumcision'
    p.level = 2

    # Add ABC approach icons
    try:
        left = Inches(1)
        top = Inches(4)
        height = Inches(3)
        slide.shapes.add_picture('abc_approach_icons.png', left, top, height=height)
    except:
        print("Could not add ABC approach icons")

    # Slide 15: Prevention and Control - Secondary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention and Control - Secondary'
    tf = body_shape.text_frame
    tf.text = 'Secondary Prevention'

    p = tf.add_paragraph()
    p.text = 'Screening Programs:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Regular STI screening for high-risk groups'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Antenatal screening (syphilis, HIV)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Targeted interventions (TI) for high-risk populations'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'NACO Programs:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Integrated Counseling and Testing Centers (ICTCs)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Prevention of Parent-to-Child Transmission (PPTCT)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Link Worker Scheme for contact tracing'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Key Achievements:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 1,381 ICTCs across India'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 95% antenatal coverage for HIV testing'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 800 million condoms distributed annually'
    p.level = 2

    # Slide 16: Prevention and Control - Tertiary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention and Control - Tertiary'
    tf = body_shape.text_frame
    tf.text = 'Tertiary Prevention'

    p = tf.add_paragraph()
    p.text = 'Management of Complications:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ART centers for HIV care'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• STD clinics for follow-up'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Community Care Centers (CCCs)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Support Services:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Positive People Networks'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Counseling and psychosocial support'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Rehabilitation programs'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Surveillance:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HIV Sentinel Surveillance'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Integrated Disease Surveillance Program'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Regular reporting and monitoring'
    p.level = 2

    # Slide 17: Challenges in India
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Challenges in India'
    tf = body_shape.text_frame
    tf.text = 'Barriers to Effective Control'

    p = tf.add_paragraph()
    p.text = 'Social and Cultural:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Stigma and discrimination'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Gender inequalities'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Religious and caste factors'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Limited sexuality education'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Healthcare System:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Rural-urban disparities'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Shortage of trained providers'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Drug stockouts'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Weak surveillance systems'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Behavioral Factors:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Low condom use'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Multiple concurrent partnerships'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Alcohol and drug use'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Migration and mobility'
    p.level = 2

    # Slide 18: Future Directions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Future Directions'
    tf = body_shape.text_frame
    tf.text = 'Way Forward'

    p = tf.add_paragraph()
    p.text = 'Strengthening Programs:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Comprehensive sexuality education in schools'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Integration of STI services with primary healthcare'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Task shifting to nurses and community health workers'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Digital health solutions for follow-up'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Research Priorities:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Vaccine development'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Point-of-care diagnostics'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Drug resistance surveillance'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Behavioral interventions'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Global Targets:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 90% reduction in syphilis incidence by 2030'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Elimination of MTCT of HIV and syphilis'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Improved access to STI services'
    p.level = 2

    # Slide 19: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Key Takeaways'
    tf = body_shape.text_frame
    tf.text = 'Summary'

    p = tf.add_paragraph()
    p.text = '1. STDs are major public health problem with significant burden in India'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '2. Most are asymptomatic, requiring active screening'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '3. Syndromic management and partner treatment are key'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '4. Prevention through ABC approach, vaccines, and condoms'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '5. NACO programs provide framework for comprehensive control'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '6. Cultural sensitivity and community engagement essential'
    p.level = 1

    # Slide 20: Q&A Session
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Q&A Session'
    tf = body_shape.text_frame
    tf.text = 'Questions and Discussion'

    p = tf.add_paragraph()
    p.text = 'Thank you for your attention!'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'References:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• NACO STI Management Guidelines (2020)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• WHO Guidelines for STI Management (2016)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• CDC STD Treatment Guidelines (2021)'
    p.level = 2

    prs.save('STD_One_Hour_Class_With_Images.pptx')
    print("STD presentation with images created successfully!")

if __name__ == "__main__":
    create_std_presentation_with_images()
