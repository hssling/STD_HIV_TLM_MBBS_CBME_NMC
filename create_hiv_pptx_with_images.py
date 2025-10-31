from pptx import Presentation
from pptx.util import Inches

def create_hiv_presentation_with_images():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "HIV/AIDS: Comprehensive Management in Indian Context"
    subtitle.text = "MBBS 3rd Year Class\nDuration: 60 minutes\nAuthor: Dr Siddalingaiah H S, Professor, Community Medicine, SIMSRH, Tumkur\nEmail: hssling@yahoo.com | Phone: 8941087719\nDate: [Date]"

    # Slide 2: What is HIV/AIDS?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'What is HIV/AIDS?'
    tf = body_shape.text_frame
    tf.text = 'Definition and Overview'

    p = tf.add_paragraph()
    p.text = 'Definition:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HIV: Human Immunodeficiency Virus'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• AIDS: Acquired Immune Deficiency Syndrome'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Retrovirus that attacks CD4+ T lymphocytes'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Leads to progressive immunosuppression'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Key Facts:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Discovered in 1983'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 39 million people living with HIV globally (2022)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 23.1 lakh PLHIV in India (NACO 2023)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Chronic manageable condition with ART'
    p.level = 2

    # Slide 3: HIV Virology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'HIV Virology'
    tf = body_shape.text_frame
    tf.text = 'Virus Structure and Replication'

    p = tf.add_paragraph()
    p.text = 'Structure:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Envelope: GP120 and GP41 proteins'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Core: Capsid containing RNA genome'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Enzymes: Reverse transcriptase, integrase, protease'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Receptors: CD4, CCR5/CXCR4 co-receptors'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Replication Cycle:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '1. Attachment: GP120 binds to CD4 receptor'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '2. Entry: Fusion with host cell membrane'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '3. Reverse Transcription: RNA → DNA'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '4. Integration: Viral DNA into host genome'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '5. Transcription: Viral mRNA production'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '6. Assembly & Budding: New virions released'
    p.level = 2

    # Slide 4: Pathogenesis and Natural History
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Pathogenesis and Natural History'
    tf = body_shape.text_frame
    tf.text = 'Disease Progression'

    # Add HIV progression timeline
    try:
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5)
        slide.shapes.add_picture('hiv_progression_timeline.png', left, top, height=height)
    except:
        print("Could not add HIV progression timeline")

    # Slide 5: Epidemiology - Global Burden
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Epidemiology - Global Burden'
    tf = body_shape.text_frame
    tf.text = 'Global HIV Statistics'

    p = tf.add_paragraph()
    p.text = 'UNAIDS 2023 Report:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 39 million people living with HIV'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 1.3 million new infections annually'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 630,000 AIDS-related deaths'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 29.8 million on antiretroviral therapy'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Regional Distribution:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Sub-Saharan Africa: 25.7 million (66% of global total)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Asia-Pacific: 5.9 million'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Western & Central Europe/North America: 2.2 million'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Eastern Europe & Central Asia: 1.5 million'
    p.level = 2

    # Slide 6: Epidemiology - Indian Context
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Epidemiology - Indian Context'
    tf = body_shape.text_frame
    tf.text = 'HIV in India: Progress and Challenges'

    p = tf.add_paragraph()
    p.text = 'Current Status (NACO 2023):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Adult prevalence: 0.22%'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• People living with HIV: 23.1 lakh'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• New infections: ~58,000 annually'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• AIDS-related deaths: ~15,000 annually'
    p.level = 2

    # Add Indian HIV transmission pie chart
    try:
        left = Inches(4)
        top = Inches(2)
        height = Inches(4)
        slide.shapes.add_picture('indian_hiv_transmission_pie.png', left, top, height=height)
    except:
        print("Could not add Indian HIV transmission pie chart")

    # Slide 7: Transmission Routes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Transmission Routes'
    tf = body_shape.text_frame
    tf.text = 'How HIV Spreads'

    p = tf.add_paragraph()
    p.text = 'Sexual Transmission (Primary Route):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Unprotected vaginal intercourse'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Unprotected anal intercourse'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Oral sex (less efficient)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Multiple partners increase risk'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Blood-Borne Transmission:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Sharing contaminated needles/syringes'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Blood transfusions (rare in screened blood)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Organ transplantation'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Mother-to-child transmission'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Risk Factors for Transmission:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• High viral load (acute infection, untreated)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Co-infections (STDs increase risk)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Lack of circumcision (male)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Traumatic sex, bleeding'
    p.level = 2

    # Slide 8: Clinical Features - Stages
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Clinical Features - Stages'
    tf = body_shape.text_frame
    tf.text = 'Clinical Presentation'

    p = tf.add_paragraph()
    p.text = 'Acute HIV Infection (Seroconversion):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Fever, rash, fatigue (70-90%)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Sore throat, lymphadenopathy'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Myalgia, arthralgia'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Headache, nausea'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Duration: 1-4 weeks'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Asymptomatic Stage:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• No symptoms for 8-10 years'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Gradual CD4 decline'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Persistent lymphadenopathy possible'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Symptomatic HIV:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Unexplained weight loss (>10%)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Chronic diarrhea'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Persistent fever'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Oral candidiasis'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Herpes zoster'
    p.level = 2

    # Slide 9: Diagnosis - Testing Algorithm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Diagnosis - Testing Algorithm'
    tf = body_shape.text_frame
    tf.text = 'HIV Diagnostic Strategy'

    p = tf.add_paragraph()
    p.text = 'NACO Testing Algorithm (2023):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Step 1: Screening Tests'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• ELISA/ECLIA (Enzyme/Chemiluminescent Immunoassay)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Rapid tests (fingerstick/oral fluid)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Combination assays (Ab + p24 Ag)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = 'Step 2: Confirmatory Tests'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Western blot (gold standard)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• Line immunoassay (cheaper alternative)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = '• HIV-1 RNA PCR (for infants <18 months)'
    p.level = 3

    p = tf.add_paragraph()
    p.text = 'Step 3: Tie-breaker (if discordant)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Different assay or HIV-1 RNA test'
    p.level = 3

    p = tf.add_paragraph()
    p.text = 'Window Periods:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Antibody tests: 4-12 weeks'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Combination tests: 2-4 weeks'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• RNA PCR: 10-14 days'
    p.level = 2

    # Slide 10: Laboratory Monitoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Laboratory Monitoring'
    tf = body_shape.text_frame
    tf.text = 'Disease Monitoring Parameters'

    # Add CD4 monitoring chart
    try:
        left = Inches(0.5)
        top = Inches(2)
        height = Inches(4.5)
        slide.shapes.add_picture('cd4_monitoring_chart.png', left, top, height=height)
    except:
        print("Could not add CD4 monitoring chart")

    # Slide 11: Antiretroviral Therapy (ART)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Antiretroviral Therapy (ART)'
    tf = body_shape.text_frame
    tf.text = 'Treatment Principles'

    p = tf.add_paragraph()
    p.text = 'NACO ART Guidelines 2023:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Test and Treat: ART for all PLHIV regardless of CD4'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Preferred First-Line: TLD (Tenofovir + Lamivudine + Dolutegravir)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Alternative First-Line: TLE (Tenofovir + Lamivudine + Efavirenz)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Second-Line: Protease inhibitor-based regimens'
    p.level = 2

    # Add ART regimen comparison
    try:
        left = Inches(0.5)
        top = Inches(3.5)
        height = Inches(3)
        slide.shapes.add_picture('art_regimen_comparison.png', left, top, height=height)
    except:
        print("Could not add ART regimen comparison")

    # Slide 12: ART Management and Side Effects
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'ART Management and Side Effects'
    tf = body_shape.text_frame
    tf.text = 'Treatment Monitoring and Challenges'

    p = tf.add_paragraph()
    p.text = 'Adherence Strategies:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Fixed-dose combinations'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Once-daily regimens'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Pill organizers, reminders'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Linkage to daily routines'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Common Side Effects:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Efavirenz: CNS effects (dizziness, nightmares)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Tenofovir: Renal toxicity, Fanconi syndrome'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Zidovudine: Anemia, neutropenia'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Dolutegravir: Minimal side effects'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Drug Interactions:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Rifampicin reduces ART levels'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Hormonal contraceptives'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Traditional medicines'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Recreational drugs'
    p.level = 2

    # Slide 13: Opportunistic Infections
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Opportunistic Infections'
    tf = body_shape.text_frame
    tf.text = 'OI Prevention and Management'

    p = tf.add_paragraph()
    p.text = 'Common OIs in India:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Tuberculosis: Most common OI (10% of HIV patients)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Pneumocystis pneumonia: CD4 <200'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Toxoplasma encephalitis: CD4 <100'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Cryptococcal meningitis: CD4 <100'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Candidiasis: Oral/esophageal'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Prophylaxis Guidelines:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Cotrimoxazole: CD4 <200 (prevents PCP, Toxoplasma, bacterial infections)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• INH: For latent TB (300mg daily × 6-9 months)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Azithromycin: For MAC prevention (CD4 <50)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Key Points:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Early ART prevents most OIs'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Prophylaxis reduces morbidity/mortality'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• TB-HIV co-infection common in India'
    p.level = 2

    # Slide 14: Prevention Strategies - Primary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention Strategies - Primary'
    tf = body_shape.text_frame
    tf.text = 'HIV Prevention Approaches'

    # Add U=U symbol
    try:
        left = Inches(4)
        top = Inches(1.5)
        height = Inches(3)
        slide.shapes.add_picture('u_equals_u_symbol.png', left, top, height=height)
    except:
        print("Could not add U=U symbol")

    p = tf.add_paragraph()
    p.text = 'Treatment as Prevention (TasP):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• U=U: Undetectable = Untransmittable'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Viral suppression prevents sexual transmission'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 96% reduction in transmission risk'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Pre-Exposure Prophylaxis (PrEP):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Daily TDF/FTC for high-risk individuals'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 99% effective when adherent'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Available through NACO since 2017'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Post-Exposure Prophylaxis (PEP):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 28-day ART regimen within 72 hours'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• For occupational/non-occupational exposure'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Emergency prevention strategy'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Biomedical Prevention:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Male circumcision (60% risk reduction)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Vaccines (in development)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Microbicides (research stage)'
    p.level = 2

    # Slide 15: Prevention Strategies - Secondary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention Strategies - Secondary'
    tf = body_shape.text_frame
    tf.text = 'Prevention of Transmission'

    p = tf.add_paragraph()
    p.text = 'ABC Approach:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Abstain from sex'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Be faithful to one uninfected partner'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Condoms consistently and correctly'
    p.level = 2

    # Add ABC approach icons
    try:
        left = Inches(4)
        top = Inches(2)
        height = Inches(3)
        slide.shapes.add_picture('abc_approach_icons.png', left, top, height=height)
    except:
        print("Could not add ABC approach icons")

    p = tf.add_paragraph()
    p.text = 'NACO Prevention Programs:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Targeted Interventions: For high-risk groups (MSM, FSWs, IDUs)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Condom Promotion: 800 million condoms distributed annually'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Blood Safety: 100% voluntary blood donation'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• ICTCs: 1,381 centers for testing and counseling'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Prevention of Mother-to-Child Transmission (PMTCT):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Antenatal HIV testing (95% coverage)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• ART for pregnant women'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Safe delivery practices'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Infant prophylaxis and testing'
    p.level = 2

    # Slide 16: Prevention Strategies - Tertiary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Prevention Strategies - Tertiary'
    tf = body_shape.text_frame
    tf.text = 'Comprehensive Care and Support'

    p = tf.add_paragraph()
    p.text = 'ART Centers and Link ART Centers:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Decentralized ART delivery'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 1,200+ centers across India'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Free lifelong treatment'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Community Care Centers (CCCs):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Nutritional support'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Psychosocial counseling'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Adherence support'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Support Networks:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Positive People Networks'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• PLHIV groups'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Peer educators'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Mental health services'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Surveillance and Monitoring:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• HIV Sentinel Surveillance'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Case reporting'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Program evaluation'
    p.level = 2

    # Slide 17: Challenges in India
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Challenges in India'
    tf = body_shape.text_frame
    tf.text = 'Barriers to HIV Control'

    p = tf.add_paragraph()
    p.text = 'Social and Cultural:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Stigma and discrimination'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Gender inequalities'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Marginalized communities (MSM, transgender)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Limited sexuality education'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Healthcare System:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Rural-urban disparities'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Human resource shortages'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Drug stockouts'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Weak referral systems'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Behavioral Factors:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Migration and mobility'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Concurrent sexual partnerships'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Alcohol and drug use'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Non-adherence to treatment'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Biological Factors:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Co-infections (TB, viral hepatitis)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Drug resistance emergence'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Late presentation for care'
    p.level = 2

    # Slide 18: Success Story and Future
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Success Story and Future'
    tf = body_shape.text_frame
    tf.text = "India's HIV Response"

    p = tf.add_paragraph()
    p.text = 'Achievements (2007-2023):'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 66% reduction in new infections'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 80% ART coverage'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 90% viral suppression rates'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• 15,000 AIDS deaths annually (down from 100,000+)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'NACP Phases:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Phase I-IV: Building infrastructure'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Phase V (2017-2021): Test and treat'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Phase VI (2021-2026): Ending AIDS by 2030'
    p.level = 2

    # Add prevention pyramid
    try:
        left = Inches(0.5)
        top = Inches(4)
        height = Inches(3)
        slide.shapes.add_picture('prevention_pyramid.png', left, top, height=height)
    except:
        print("Could not add prevention pyramid")

    p = tf.add_paragraph()
    p.text = 'Future Directions:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 95-95-95 targets by 2030'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Integration with general healthcare'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Digital solutions for follow-up'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Research in vaccines and cure'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Global Goals:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• End AIDS epidemic by 2030'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• <200,000 new infections annually'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• Zero discrimination'
    p.level = 2

    # Slide 19: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Key Takeaways'
    tf = body_shape.text_frame
    tf.text = 'Summary'

    p = tf.add_paragraph()
    p.text = '1. HIV is a manageable chronic condition with ART'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '2. Early diagnosis and treatment prevent complications'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '3. U=U: Undetectable viral load prevents transmission'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '4. Prevention through TasP, PrEP, PEP, and condoms'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '5. India\'s response shows what commitment can achieve'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '6. Stigma reduction and community engagement essential'
    p.level = 2

    # Slide 20: Q&A Session
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Q&A Session'
    tf = body_shape.text_frame
    tf.text = 'Questions and Discussion'

    p = tf.add_paragraph()
    p.text = 'Thank you for your attention!'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'References:'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• NACO ART Guidelines 2023'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• WHO HIV Guidelines 2021'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '• UNAIDS Global AIDS Update 2023'
    p.level = 2

    prs.save('HIV_One_Hour_Class_With_Images.pptx')
    print("HIV presentation with images created successfully!")

if __name__ == "__main__":
    create_hiv_presentation_with_images()
